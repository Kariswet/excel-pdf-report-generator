from loguru import logger
from pptx.chart.data import CategoryChartData
from datetime import datetime
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

class ChartAndOthers:
    def __init__(self):
        pass

    def _replace_chart(self, slide, C_index, C_data):
        chart = 0
        for shape in slide.shapes:
            if shape.has_chart:
                chart += 1
                if chart == C_index + 1:
                    chart = shape.chart
                    chart.replace_data(C_data)
                    logger.debug(f"Chart index {C_index} replaced")
                    return
        logger.debug(f"Chart index {C_index} not found")

    def generate_line_chart(self,presentation, data, s_idx, c_idx, type):
        data_chart = data

        if not data_chart:
            categories = []
            series_data = {}
        else:
            categories = []
            series_data = {}

        platforms = set()

        for item in data_chart:
            for platform in item.get("inner_buckets", []):
                platforms.add(platform["key"])

        platforms = list(platforms)

        for item in data_chart:
            # categories.append(item["key_as_string"])
            date_only = item["key_as_string"][:10]  # "2025-07-01"
            categories.append(date_only)

            daily_data = {platform: 0 for platform in platforms}

            for platform in item.get("inner_buckets", []):
                daily_data[platform["key"]] = platform["doc_count"]
                if type == "engagement":
                    daily_data[platform["key"]] = platform["value"]

            for platform in platforms:
                if platform not in series_data:
                    series_data[platform] = []
                series_data[platform].append(daily_data[platform])

        chart_data = CategoryChartData()
        if categories:
            chart_data.categories = categories
            for platform, values in series_data.items():
                chart_data.add_series(platform, values)
        else:
            chart_data.categories = [""]
            chart_data.add_series("", [0])

        slide_index = s_idx
        chart_index_to_replace = c_idx

        for i, slide in enumerate(presentation.slides):
            if i == slide_index:
                self._replace_chart(slide, chart_index_to_replace, chart_data)
                break
    
    def generate_chart_string_key(self, presentation, data, s_idx, c_idx, type):
        data_chart = data
        if not data_chart:
            categories = []
            values = []
        else:
            categories = []
            values = []

        sorted_data = sorted(data_chart, key=lambda x: x['doc_count'], reverse=False)
        for item in sorted_data:
            categories.append(item['key'])
            if type == "engagement":
                values.append(f"{item['value']:.2f}")
            else:
                values.append(item['doc_count'])
        
        chart_data = CategoryChartData()
        if categories:
            chart_data.categories = categories
            chart_data.add_series("", values)
        else:
            chart_data.categories = [""]
            chart_data.add_series("", [0])
        
        slide_index = s_idx
        chart_index_to_replace = c_idx
        
        for i, slide in enumerate(presentation.slides):
            if i == slide_index:
                self._replace_chart(slide, chart_index_to_replace, chart_data)
                break
    
    def retrieve_data_object(self, data, key=None, sub_key=None, default_fill='', length_min=10):
        if isinstance(data, dict) and key:
            items = data.get(key, [])
        elif isinstance(data, list):
            items = data
        else:
            return [default_fill] * length_min

        values = []

        if isinstance(items, list):
            for item in items:
                if isinstance(item, dict):
                    if sub_key: 
                        value = item.get(sub_key, default_fill)
                    else:
                        value = item
                    values.append(value if value else default_fill)
                elif isinstance(item, (str, int)):
                    values.append(str(item))
                elif isinstance(item, list):
                    values.extend(item if item else [default_fill])
                else:
                    values.append(default_fill)
        else:
            values = [default_fill]

        flat_values = []
        for val in values:
            if isinstance(val, list):
                flat_values.extend(val)
            else:
                flat_values.append(val)

        while len(flat_values) < length_min:
            flat_values.append(default_fill)

        return flat_values
    
    def rank_sum(self, data_count, idx):
        data_rank = 0
        for avg in data_count:
            data_rank += avg
            devide = data_rank / len(data_count)
        data = data_count[idx] - devide
        return f"{round(data, 2)}"
    
    def generate_bar_chart(self, presentation, data, s_idx, c_idx):
        data_chart = data
        if not data_chart:
            categories = []
            values = []
        else:
            categories = []
            values = []

        sorted_data = sorted(data_chart, key=lambda x: x['key'], reverse=False)
        for item in sorted_data:
            dt = datetime.fromtimestamp(item['key'] / 1000).strftime("%d-%m-%Y")
            categories.append(dt)
            values.append(item['doc_count'])
        
        chart_data = CategoryChartData()
        if categories:
            chart_data.categories = categories
            chart_data.add_series("", values)
        else:
            chart_data.categories = [""]
            chart_data.add_series("", [0])
        
        slide_index = s_idx
        chart_index_to_replace = c_idx
        
        for i, slide in enumerate(presentation.slides):
            if i == slide_index:
                self._replace_chart(slide, chart_index_to_replace, chart_data)
                break

    def create_table(self,slide, data, start_no, left, top, width, height, col_widths, header_size):
        """Helper function to create a table on the slide with adjustable column widths."""
        if data.empty:
            return  # Avoid creating an empty table
        
        data.insert(0, "No", range(start_no, start_no + len(data)))

        row, cols = data.shape
        table = slide.shapes.add_table(row + 1, cols, left, top, width, height).table

        # Adjust column widths (if provided)
        if col_widths:
            for col_idx, col_width in enumerate(col_widths):
                table.columns[col_idx].width = Cm(col_width)

        # Set small row height (PowerPoint still auto-adjusts)
        for row_idx in range(row + 1):
            table.rows[row_idx].height = Cm(0.75) 

        # Apply header styling
        for col_idx, col_name in enumerate(data.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(header_size)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(200, 100, 50)  # Header background color
            cell.fill.fore_color.transparency = 1.0
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

        # Fill table with data
        for row_idx, row in data.iterrows():
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(8)  # Font size
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)   # Black text   
    
    def generate_table(self, prs, s_idx, data, case):
        """Generates tables on a PowerPoint slide with at most 17 rows each."""
        slide = prs.slides[s_idx]

        # Define column widths (adjust as needed)
        if case == "kpi":
            column_widths = [1.1, 6, 1.4, 1.4, 1.4, 1.4] 
            left = Cm(3.81)
            left2 = Cm(18.01)
            top = Cm(4.4)
            width = Cm(13)
            height = Cm(14.4)
            header_size = 11
            
        elif case == "dpi":
            column_widths = [1.1, 9, 1.2]
            left = Cm(4)
            left2 = Cm(18.01)
            top = Cm(4.4)
            width = Cm(14.4)
            height = Cm(14.4)
            header_size= 11
            
        elif case == "dii":
            column_widths = [1.1, 9, 1.2]
            left = Cm(4)
            left2 = Cm(18.01)
            top = Cm(4.4)
            width = Cm(14.4)
            height = Cm(14.4)
            header_size= 11
            
        elif case == "ppi":  
            column_widths = [1.1, 9, 1.2]
            left = Cm(4)
            left2 = Cm(18.01)
            top = Cm(4.4)
            width = Cm(14.4)
            height = Cm(14.4)
            header_size= 11
            

        # First table (Rows 1-17)
        data1 = data.iloc[:17].copy()
        self.create_table(slide, data1, start_no=1, left=left, top=top, 
                    width=width, height=height, col_widths=column_widths, header_size=header_size)

        # Second table (Rows 18-34) if applicable
        data2 = data.iloc[17:].reset_index(drop=True).copy()
        if not data2.empty:
            self.create_table(slide, data2, start_no=18, left=left2, top=top, 
                    width=width, height=height, col_widths=column_widths, header_size=header_size)
            
        logger.debug(f"Table at slide index {s_idx} created")